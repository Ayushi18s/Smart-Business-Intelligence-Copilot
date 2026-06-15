from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation


# =========================
# PDF REPORT GENERATOR
# =========================
def generate_pdf_report(text, filename="report.pdf"):
    doc = SimpleDocTemplate(filename)
    styles = getSampleStyleSheet()

    content = []

    for line in text.split("\n"):
        content.append(Paragraph(line, styles["Normal"]))
        content.append(Spacer(1, 8))

    doc.build(content)
    return filename


# =========================
# PPT GENERATOR
# =========================
def generate_ppt_report(text, filename="report.pptx"):
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = "AI Business Report"
    body.text = text[:2000]  # limit for slide

    prs.save(filename)
    return filename